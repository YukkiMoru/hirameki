import os
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image  # 追加

base_dir = "hirameki_files"
prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]

for folder in sorted(os.listdir(base_dir)):
    folder_path = os.path.join(base_dir, folder)
    if not os.path.isdir(folder_path):
        continue

    slide = prs.slides.add_slide(blank_slide_layout)
    left = top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.7)
    title_shape = slide.shapes.add_textbox(left, top, width, height)
    title_shape.text = folder
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)

    images = [f for f in os.listdir(folder_path) if f.lower().endswith(('.png', '.jpg', '.jpeg'))]
    images.sort()
    cols = 4
    max_img_size = Inches(2.0)  # 最大サイズを小さく
    spacing = Inches(0.3)
    for idx, img_name in enumerate(images):
        row = idx // cols
        col = idx % cols
        img_path = os.path.join(folder_path, img_name)
        # 画像サイズ取得して縦横比を維持
        with Image.open(img_path) as im:
            w, h = im.size
            ratio = min(max_img_size / w, max_img_size / h)
            disp_w = w * ratio
            disp_h = h * ratio
        left = Inches(0.5) + col * (max_img_size + spacing)
        top = Inches(1.2) + row * (max_img_size + spacing)
        slide.shapes.add_picture(img_path, left, top, width=disp_w, height=disp_h)

prs.save("hirameki_colors.pptx")
print("PowerPointファイルを作成しました: hirameki_colors.pptx")